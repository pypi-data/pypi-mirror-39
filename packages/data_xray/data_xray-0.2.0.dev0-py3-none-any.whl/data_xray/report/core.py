from data_xray.modules import *
from data_xray.scan import PlotImage
#some utilies to work with powerpoint


class SummaryPPT(object):
    def __init__(self, pname=None, new=False, fdict=None, **kwargs):

        if fdict is None:
            print('please specify data to summarize')
            return
        else:
            self.fdict = fdict

        self.topdir = os.path.commonpath([j.ds.fname for j in fdict]) + '/'
        self.presentation_name = pname
        self.pptx_file_name = self.topdir + self.presentation_name + '.pptx'

        self.new = new
        self.init_ppt()
        self.insert_images()

        try:
            self.pres.save(self.pptx.file_name)
            print('images stored in : ' + self.pptx_file_name )
        except:
            print('something wrong with saving the presentation file')

    def init_ppt(self):
        """
        Initialize powerpoint presentation
        :param pname:
        :param new:
        :return:
        """
        if self.presentation_name is None:
            newpres = Presentation()
            newpres.notes_master.name = 'sum1.pptx'
            newpres.save(newpres.notes_master.name)
            self.pres = newpres
        elif self.new:
            newpres = Presentation()
            newpres.notes_master.name = self.presentation_name
            newpres.save(newpres.notes_master.name)
            self.pres = newpres
        else:
            pres = Presentation(self.presentation_name)
            pres.notes_master.name = self.presentation_name
            self.pres = pres

    def insert_images(self, chanselect='Z'):
        """
        Dumpt a batch of images into a powerpoint

        :param chanselect:
        :param fdict:
        :param topdir:
        :return:
        """
        #for folder in (self.fdict.keys()):
        #        self.text_to_slide(folder)

        for fj in self.fdict:
            # TextToSlide(fj.fname,pres=pres)
            try:
                f3, a3 = plt.subplots(1, 1)
                d2 = PlotImage(fj, chan=chanselect, ax=a3, high_pass=None)
                self.fig_to_ppt([f3], leftop=[3, 2], txt=fj.fname)
                print(os.path.basename(fj.ds.fname) + ' imported')

            except:
                print(os.path.basename(fj.ds.fname) + ' failed')

    def png_to_ppt(self, pngfile, ttl = []):
       """
       Plop a PNG file into powerpoint slide
       :param pngfile:
       :param pres:
       :param ttl:
       :return:
       """

       #blank_slide_layout = pres.slide_layouts[6]
       title_slide_layout = self.pres.slide_layouts[9]

       left = top = Inches(1)

       slide = self.pres.slides.add_slide(title_slide_layout)
       slide.shapes.add_picture(pngfile, left, top)
       subtitle = slide.placeholders[1]
       title = slide.shapes.title
       if len(ttl):
           subtitle.text = ttl


    def fig_to_ppt(self, figs, leftop=[0,1.5], txt=None):
        """
        Plop figures into powerpoint
        :param figs:
        :param pres:
        :param leftop:
        :param txt:
        :return:
        """
        #savepptx needs to be a full path. If None is provided the default presentation
        #will be created with a name sum1.pptx in the current folder
        from pptx.util import Inches

        blank_slide_layout = self.pres.slide_layouts[5]
        left = Inches(leftop[0])
        top = Inches(leftop[1])

        tmp_path = 't1.png'
        for figp in figs:
            plt.savefig(tmp_path, transparent=1, format='png', dpi=300, bbox_inches = 'tight')
            slide = self.pres.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(tmp_path, left, top)

        if txt is not None:
            self.text_to_slide(txt, slide=slide)

    def text_to_slide(self, txt, slide=None):
        """
        convert text to slide

        :param txt:
        :param pres:
        :param slide:
        :return:
        """
        from pptx.util import Pt
        #title = slide.shapes.title
        #subtitle = slide.placeholders[1]

       # title.text = "Hello, World!"
        #subtitle.text = "python-pptx was here!"

       # prs.save('test.pptx')

        from pptx.util import Inches

        if self.pres == None:
            print('please init presentation')
        else:
            if slide is None:
                bullet_slide_layout = self.pres.slide_layouts[5]
                slide = self.pres.slides.add_slide(bullet_slide_layout)

            shapes = slide.shapes

            countshapes = 0

            #just catch the first shape object with a frame in it
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elif countshapes > 0:
                    tframe = shape.text_frame
                    tframe.clear()
                    #print('caught one')
                else:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    countshapes = 1

            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = txt

            font = run.font
            font.name = 'Calibri'
            font.size = Pt(12)


