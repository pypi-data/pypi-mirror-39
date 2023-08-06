from ..modules import *

#some utilies to work with powerpoint

def InitPPT(pname=None, new=False):
    if pname is None:
        newpres = Presentation()
        newpres.notes_master.name = 'sum1.pptx'
        newpres.save(newpres.notes_master.name)
        return newpres
    elif new:
        newpres = Presentation()
        newpres.notes_master.name = pname
        newpres.save(newpres.notes_master.name)
        return newpres
    else:
        pres = Presentation(pname)
        pres.notes_master.name = pname
        return pres

def PNGToPPT(pngfile, pres=None, ttl = []):

   #blank_slide_layout = pres.slide_layouts[6]
   title_slide_layout = pres.slide_layouts[9]

   left = top = Inches(1)

   slide = pres.slides.add_slide(title_slide_layout)
   slide.shapes.add_picture(pngfile, left, top)
   subtitle = slide.placeholders[1]
   title = slide.shapes.title
   if len(ttl):
       subtitle.text = ttl


def FigToPPT(figs, pres=None, leftop=[0,1.5], txt=None):
    #savepptx needs to be a full path. If None is provided the default presentation
    #will be created with a name sum1.pptx in the current folder
    from pptx.util import Inches

    if pres == None:
        pres = Presentation()

    blank_slide_layout = pres.slide_layouts[5]
    left = Inches(leftop[0])
    top = Inches(leftop[1])

    tmp_path = 't1.png'
    for figp in figs:
        plt.savefig(tmp_path, transparent=1, format='png', dpi=300, bbox_inches = 'tight')
        slide = pres.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(tmp_path, left, top)

    if txt is not None:
        TextToSlide(txt, pres=pres, slide=slide)

    pptxfile = 'sum1.pptx' if pres.notes_master.name == '' else pres.notes_master.name
    #pres.save(pptxfile)


def TextToSlide(txt, pres=None, slide=None):
    from pptx.util import Pt
    #title = slide.shapes.title
    #subtitle = slide.placeholders[1]

   # title.text = "Hello, World!"
    #subtitle.text = "python-pptx was here!"

   # prs.save('test.pptx')

    from pptx.util import Inches

    if pres == None:
        print('please init presentation')
    else:
        if slide is None:
            bullet_slide_layout = pres.slide_layouts[5]
            slide = pres.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes

        countshapes = 0

        #just catch the first shape object with a frame in it
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            elif countshapes > 0:
                tframe = shape.text_frame
                tframe.clear()
                #print('caught one')
            else:
                text_frame = shape.text_frame
                text_frame.clear()
                countshapes = 1

        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = txt

        font = run.font
        font.name = 'Calibri'
        font.size = Pt(12)


        #title_shape = shapes.title
        #body_shape = shapes.placeholders[1]

        #title_shape.text = 'Adding a Bullet Slide'

        #tf = body_shape.text_frame
        #tf.text = txt

        #p = tf.add_paragraph()
        #p.text = 'Use _TextFrame.text for first bullet'
        #p.level = 1

        #p = tf.add_paragraph()
        #p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
        #p.level = 2

        pptxfile = 'sum1.pptx' if pres.notes_master.name == '' else pres.notes_master.name
        #pres.save(pptxfile)